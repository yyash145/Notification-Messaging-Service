from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

image_files = [
    "./slide-1.png",
    "./slide-2.png",
    "./slide-3.png",
    "./slide-4.png",
    "./slide-5.png",
    "./slide-6.png",
    "./slide-7.png",
    "./slide-8.png",
    "./slide-9.png",
    "./slide-10.png"
]

for img in image_files:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide
    slide.shapes.add_picture(img, Inches(0), Inches(0),
                              width=prs.slide_width,
                              height=prs.slide_height)

prs.save("notifications.pptx")
